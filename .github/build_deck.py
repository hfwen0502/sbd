"""Build the SBD Python Walkthrough deck as a PPTX.

Concise, draft-quality: trims the .md commentary down to slide-scale,
highlights only the critical lines in code blocks. Iterate on this,
not on the .pptx directly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# -----------------------------------------------------------------------------
# Style constants
# -----------------------------------------------------------------------------
PRIMARY  = RGBColor(0x0F, 0x62, 0xFE)   # IBM blue
ACCENT   = RGBColor(0xDA, 0x1E, 0x28)   # IBM red
DIM      = RGBColor(0x69, 0x76, 0x7A)   # gray for muted code lines
DARK     = RGBColor(0x16, 0x16, 0x16)
LIGHT    = RGBColor(0xF4, 0xF4, 0xF4)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
CODE_FONT  = "Menlo"   # macOS-friendly monospace; Consolas works too

# Slide layout (16:9)
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------
def add_title(slide, title, subtitle=None):
    """Add a title bar at the top of a blank slide."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                    Inches(12.3), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    r = p.runs[0]
    r.font.name = TITLE_FONT
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        r2 = p2.runs[0]
        r2.font.name = TITLE_FONT
        r2.font.size = Pt(16)
        r2.font.color.rgb = PRIMARY
    return box


def add_text_box(slide, left, top, width, height, lines,
                 font=BODY_FONT, size=Pt(16), color=DARK,
                 bold_first=False):
    """Add a textbox with a list of plain-text bullets / lines."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line else " "
        r = p.runs[0]
        r.font.name = font
        r.font.size = size
        r.font.color.rgb = color
        if bold_first and i == 0:
            r.font.bold = True
    return box


def add_code_box(slide, left, top, width, height, lines,
                  highlight_idx=None, highlight_color=ACCENT,
                  bg=LIGHT, size=Pt(14)):
    """Add a code snippet as a textbox, optional line-index highlight set."""
    if highlight_idx is None:
        highlight_idx = set()
    else:
        highlight_idx = set(highlight_idx)

    # Background rectangle
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       left, top, width, height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False

    # Text on top
    box = slide.shapes.add_textbox(left + Emu(50000), top + Emu(50000),
                                    width - Emu(100000), height - Emu(100000))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line else " "
        r = p.runs[0]
        r.font.name = CODE_FONT
        r.font.size = size
        if i in highlight_idx:
            r.font.color.rgb = highlight_color
            r.font.bold = True
        else:
            r.font.color.rgb = DARK if line.strip() and not line.lstrip().startswith("#") else DIM
    return box


def add_table(slide, left, top, width, height, rows):
    """Simple table from a list of row-tuples (first row is header)."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(14)
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        run.font.color.rgb = DARK
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r_idx % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    return tbl


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05),
                                    Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.name = BODY_FONT
    r.font.size = Pt(10)
    r.font.color.rgb = DIM
    r.font.italic = True


# =============================================================================
# Slide 1 — Title
# =============================================================================
slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5),
                                Inches(12.3), Inches(2.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SBD Python Walkthrough"
r = p.runs[0]
r.font.name = TITLE_FONT
r.font.size = Pt(48)
r.font.bold = True
r.font.color.rgb = DARK

p2 = tf.add_paragraph()
p2.text = "GPU-accelerated alternative to Dice in qiskit-addon-sqd"
p2.runs[0].font.name = TITLE_FONT
p2.runs[0].font.size = Pt(24)
p2.runs[0].font.color.rgb = PRIMARY

p3 = tf.add_paragraph()
p3.text = ""
p4 = tf.add_paragraph()
p4.text = "Co-presented by [systems lead] and [chemistry lead]"
p4.runs[0].font.name = TITLE_FONT
p4.runs[0].font.size = Pt(16)
p4.runs[0].font.italic = True
p4.runs[0].font.color.rgb = DIM


# =============================================================================
# Slide 2 — SBD vs Dice (positioning)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "SBD vs Dice — same role, different envelope",
          "Both diagonalize the projected Hamiltonian inside qiskit-addon-sqd's iteration loop")

add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.6), [
    ("",                          "Dice (SHCI)",                       "SBD"),
    ("Hardware",                  "CPU + MPI",                         "CPU + MPI + GPU (Thrust / OMP-offload)"),
    ("Process model",             "subprocess + CLI + binary file I/O","in-process pybind11 module"),
    ("Determinant selection",     "internal heat-bath (eps knob)",     "external (caller-provided list)"),
    ("Orbital ceiling (data-type)","128  (16-byte determinant addr)",  "160 (default) · 512 (bit_length=64) · configurable"),
    ("Integration with sqd-addon","stock package",                     "MPI-aware fork (@patch-ferminon-sbd)"),
])

add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
             ["Where it lives"], size=Pt(16), color=PRIMARY, bold_first=True)
add_table(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.5), [
    ("Layer",                                            "Repo"),
    ("Upstream SBD (C++ core, RIKEN CCS)",               "github.com/r-ccs-cms/sbd"),
    ("Python wrapper (fork with python/ bindings)",      "github.com/hfwen0502/sbd"),
    ("Patched qiskit-addon-sqd (MPI-aware solver hook)", "github.com/hfwen0502/qiskit-addon-sqd  (branch: patch-ferminon-sbd)"),
])

add_text_box(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
             ["▶  Quick start (chemist-friendly):  python/examples/run_sqd_sbd.ipynb"],
             size=Pt(13), color=PRIMARY, bold_first=True)

add_text_box(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
             ["Same FCIDUMP input, same RDM outputs, same role in the SQD loop. Process model and hardware envelope are what change."],
             size=Pt(12), color=DIM)

add_footer(slide, "Slide 2 — positioning")


# =============================================================================
# Slide 3 — Architecture: subprocess vs in-process
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Subprocess + file I/O   vs   in-process pybind11",
          "What happens per SQD iteration")

# Two side-by-side stacked-flow boxes
LEFT = Inches(0.5);  WIDTH = Inches(6.0)
RIGHT = Inches(6.83); RWIDTH = Inches(6.0)

# Left — Dice flow
add_text_box(slide, LEFT, Inches(1.7), WIDTH, Inches(0.4),
             ["Dice path"], size=Pt(20), color=ACCENT, bold_first=True)
add_code_box(slide, LEFT, Inches(2.2), WIDTH, Inches(4.5), [
    "Python",
    "  ↓",
    "write FCIDUMP + input.dat to /tmp/dice_cli_files_*",
    "  ↓",
    "subprocess.run([\"mpirun\", ..., bin/Dice], ...)",
    "  ↓",
    "Dice CLI runs",
    "  ↓",
    "parse spin1RDM.*.txt + stdout",
    "  ↓",
    "shutil.rmtree(temp_dir)",
    "  ↓",
    "back to Python",
], highlight_idx=[2, 4, 8], size=Pt(15))

# Right — SBD flow
add_text_box(slide, RIGHT, Inches(1.7), RWIDTH, Inches(0.4),
             ["SBD path"], size=Pt(20), color=PRIMARY, bold_first=True)
add_code_box(slide, RIGHT, Inches(2.2), RWIDTH, Inches(4.5), [
    "Python",
    "  ↓",
    "sbd.tpb_diag(fcidump, adet, ...)",
    "  ↓",
    "pybind11 → C++ Davidson",
    "  ↓",
    "MPI + OpenMP + GPU (Thrust or OMP-offload)",
    "  ↓",
    "returns Python objects (energy, RDMs, wf)",
    "",
    "",
    "",
    "",
], highlight_idx=[2, 4], size=Pt(15))

# Launch command contrast at the bottom
add_text_box(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
             ["How you launch the job"], size=Pt(15), color=PRIMARY, bold_first=True)
add_code_box(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5), [
    "Dice:  python my_script.py                              # serial; Dice forks mpirun per call internally",
    "SBD:   mpirun -np N python my_script.py                 # MPI-native; you launch all ranks once",
], highlight_idx=[0, 1], size=Pt(12))

add_footer(slide, "Slide 3 — architecture contrast")


# =============================================================================
# Slide 4 — Pick your backend at runtime
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Pick your backend at runtime",
          "import sbd loads every backend that compiled — switch with one keyword arg")

add_text_box(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.4),
             ["See what's available"], size=Pt(18), color=PRIMARY, bold_first=True)
add_code_box(slide, Inches(0.5), Inches(2.05), Inches(6.0), Inches(2.0), [
    "import sbd",
    "",
    "sbd.available_backends()",
    "  → ['cpu', 'gpu', 'gpu-nvidia-omp']",
    "",
    "sbd.print_info()  # version, hardware, session",
], highlight_idx=[2, 5])

add_text_box(slide, Inches(6.83), Inches(1.6), Inches(6.0), Inches(0.4),
             ["Switch backends — three equivalent paths"],
             size=Pt(18), color=PRIMARY, bold_first=True)
add_code_box(slide, Inches(6.83), Inches(2.05), Inches(6.0), Inches(4.4), [
    "# 1. process-wide default at startup",
    "sbd.init(device='gpu')",
    "",
    "# 2. per-call override, no re-init",
    "sbd.tpb_diag(..., device='gpu-omp')",
    "",
    "# 3. through DeviceConfig (qiskit-addon-sqd path)",
    "solve_sci(..., device_config=DeviceConfig.gpu())",
    "",
    "# CLI: every example script accepts --device",
    "$ python run_sbd_diag.py --device gpu",
], highlight_idx=[1, 4, 7, 10])

add_text_box(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.0),
             ["All three return bit-equal energies.",
              "Backend choice doesn't move the eigenvalue.",
              "It moves the wallclock."],
             size=Pt(15), color=DIM)

add_footer(slide, "Slide 4 — runtime backend selection")


# =============================================================================
# Slide 5 — Under the hood: same call shape
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Same call shape, very different invocation",
          "What the user calls vs what the wrapper does internally")

# Top: user-facing API
add_text_box(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.4),
             ["What the user calls"], size=Pt(16), color=DARK, bold_first=True)
add_code_box(slide, Inches(0.5), Inches(1.95), Inches(12.3), Inches(2.0), [
    "# Dice  →  qiskit_addon_dice_solver.solve_sci   (returns SCIResult)",
    "result = solve_sci(ci_strings, one_body_tensor, two_body_tensor,",
    "                   norb=norb, nelec=nelec)",
    "energy, sci_state, occ = result.energy, result.sci_state, result.orbital_occupancies",
    "",
    "# SBD   →  sbd.sbd_solver.solve_sci             (also returns SCIResult — same API shape)",
    "result = solve_sci(ci_strings, one_body_tensor, two_body_tensor,",
    "                   norb=norb, nelec=nelec, device_config=DeviceConfig.gpu())",
    "energy, sci_state, occ = result.energy, result.sci_state, result.orbital_occupancies",
    "#                                                         ^^^^^^^^^^^^^^^^^^^^^^^^^",
    "#                                                         (avg α-occ, avg β-occ)",
], highlight_idx=[0, 5, 7])

# Bottom: side-by-side internals
LEFT = Inches(0.5);  WIDTH = Inches(6.0)
RIGHT = Inches(6.83); RWIDTH = Inches(6.0)

add_text_box(slide, LEFT, Inches(3.3), WIDTH, Inches(0.4),
             ["Inside Dice's wrapper"], size=Pt(16), color=ACCENT, bold_first=True)
add_code_box(slide, LEFT, Inches(3.7), WIDTH, Inches(3.2), [
    "dice_dir = Path(tempfile.mkdtemp(",
    "    prefix='dice_cli_files_', ...))",
    "tools.fcidump.from_integrals(",
    "    dice_dir / 'fcidump.txt', ...)",
    "_write_input_files(...)",
    "",
    "dice_call = ['mpirun', *mpirun_options,",
    "             '<…>/bin/Dice']",
    "subprocess.run(dice_call, cwd=dice_dir,",
    "               stdout=logfile, stderr=logfile)",
    "",
    "_read_dice_outputs(dice_dir, ...)",
    "shutil.rmtree(dice_dir)",
], highlight_idx=[6, 7, 8, 9], size=Pt(13))

add_text_box(slide, RIGHT, Inches(3.3), RWIDTH, Inches(0.4),
             ["Inside SBD's wrapper"], size=Pt(16), color=PRIMARY, bold_first=True)
add_code_box(slide, RIGHT, Inches(3.7), RWIDTH, Inches(3.2), [
    "def tpb_diag(fcidump, adet, bdet, sbd_data,",
    "             loadname='', savename='', device=None):",
    "    _ensure_initialized()              # MPI comm cached",
    "    backend = get_backend(device)      # cpu | gpu | gpu-omp",
    "    return backend.tpb_diag(",
    "        _global_comm,                  # MPI communicator",
    "        sbd_data, fcidump, adet, bdet,",
    "        loadname, savename,",
    "    )                                  # pybind11 → C++",
], highlight_idx=[3, 5], size=Pt(13))

add_footer(slide, "Slide 5 — under-the-hood contrast")


# =============================================================================
# Slide 6 — SQD with SBD: 4-step recipe
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "SQD with SBD — bitstrings in, energy out",
          "Same shape as plugging in qiskit-addon-dice-solver; only sci_solver= changes")

add_code_box(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3), [
    "# 1. Hamiltonian from FCIDUMP",
    "mf = tools.fcidump.to_scf('data/h2o/fcidump.txt')",
    "hcore, eri = mf.get_hcore(), ao2mo.restore(1, mf._eri, norb)",
    "",
    "# 2. Bitstrings — quantum-measurement counts → BitArray",
    "counts = json.load(open('data/h2o/count_dict.json'))",
    "bit_array = BitArray.from_bool_array(_decode_counts(counts))",
    "",
    "# 3. Wire SBD into qiskit-addon-sqd's sci_solver= slot",
    "#    (no sbd.init() / mpi_comm needed: solve_sci_batch auto-inits",
    "#     and falls back to MPI.COMM_WORLD when mpi_comm is omitted)",
    "sbd_solver = partial(",
    "    solve_sci_batch,",
    "    sbd_config = {'method': 0, 'max_it': 100, 'max_nb': 50, ...},",
    "    device_config = DeviceConfig.gpu(),       # or .cpu() / .gpu_omp()",
    "    fcidump_path = 'data/h2o/fcidump.txt',",
    ")",
    "",
    "# 4. Run the SQD self-consistent loop",
    "result = diagonalize_fermionic_hamiltonian(",
    "    hcore, eri, bit_array,",
    "    norb=norb, nelec=(num_elec_a, num_elec_b),",
    "    samples_per_batch=300, num_batches=3, max_iterations=5,",
    "    sci_solver=sbd_solver,    # ← SBD plugs in here",
    ")",
], highlight_idx=[6, 12, 23], size=Pt(13))

# Quick-start callout — point readers at the runnable notebook
add_text_box(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             ["▶  Try it interactively:  python/examples/run_sqd_sbd.ipynb     "
              "(serial Jupyter — converges to ~−76.19 Ha on h2o in ~10 s on CPU)"],
             size=Pt(13), color=PRIMARY, bold_first=True)

add_footer(slide, "Slide 6 — SQD-with-SBD recipe  ·  driver: python/examples/run_sqd_sbd.py  ·  notebook: run_sqd_sbd.ipynb")


# =============================================================================
# Slide 7 — Three backends, one decision (perf)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Three backends, one runtime decision",
          "Same source, same energies (bit-equal), different compilers — NVIDIA GPU only")

add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.6), [
    ("device=",            "Compiler",                  "When to use"),
    ("'cpu'",              "system c++",                "small problems · debugging · no GPU"),
    ("'gpu' (Thrust)",     "NVHPC nvc++",               "NVIDIA GPU · production default"),
    ("'gpu-omp' (LLVM)",   "clang++ w/ NVPTX target",   "NVIDIA GPU · alternative kernel path"),
])

add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
             ["Fe4S4 27,901 α-determinants on coreweave GB200 (post-rebuild, native sm_100):"],
             size=Pt(15), color=PRIMARY, bold_first=True)

add_table(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(1.7), [
    ("Backend",       "4×GB200 (1 node)",  "8×GB200 (2 nodes)",  "1n→2n"),
    ("Thrust",        "577 s",             "329 s",              "1.75×"),
    ("OMP-offload",   "491 s",             "304 s",              "1.62×"),
])

add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0),
             ["All energies = −326.821832430028 (bit-equal across cpu / gpu / gpu-omp and 1n / 2n).",
              "Backend choice doesn't move the eigenvalue — it moves the wallclock."],
             size=Pt(14), color=DIM)

add_footer(slide, "Slide 7 — hardware story + perf")


# =============================================================================
# Slide 8 — Roadmap (co-presenter slide)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "SBD as a GPU-accelerated SCI driver",
          "Adaptive SCI features on the singles-doubles-extend branch · co-presenter handles chemistry interpretation")

# Branch ribbon — note SBD C++ is a submodule on main, embedded on the experimental branch
add_table(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.1), [
    ("Branch",                  "SBD C++ source",                                        "What it has"),
    ("main (stable)",           "submodule → upstream r-ccs-cms/sbd (unmodified)",        "Python wrapper, three backends, SQD integration"),
    ("singles-doubles-extend",  "embedded in fork (C++ modified for new features)",       "+ variance, S+D, ERI screening, TrimSQD"),
])
add_text_box(slide, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.35),
             ["github.com/hfwen0502/sbd  ·  experimental: tree/singles-doubles-extend  ·  see VARIANCE.md"],
             size=Pt(11), color=DIM)

add_text_box(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
             ["Experimental features"],
             size=Pt(16), color=PRIMARY, bold_first=True)

add_text_box(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.9),
             ["1.  Singles + Doubles subspace expansion  (--carryover_type 4–6)",
              "      Brute-force: extend selected dets with single + same-spin double excitations.",
              "      ≡ Hamming distance ≤ 2 from each seed determinant.",
              "",
              "2.  ERI-screened S+D  (--carryover_type 7–8, --eri_threshold)  — HCI-flavor selection",
              "      Keep only excitations whose Hamiltonian coupling exceeds threshold.",
              "      Typically 20–50% of brute-force S+D, retaining the physically important excitations.",
              "",
              "3.  Variance-only mode  (--iteration 0)  +  TrimSQD (adaptive pruning)",
              "      σ² = ⟨Hψ|Hψ⟩/‖ψ‖² − E² without diagonalizing.  Two-step expand → diagonalize → variance → repeat → σ² → 0.",
             ],
             size=Pt(12))

# Concrete demo table from VARIANCE.md
add_text_box(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.4),
             ["Demo: NORB=29, nelec=(5α,5β),  seeded from 995 sampled dets"],
             size=Pt(13), color=PRIMARY, bold_first=True)
add_table(slide, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.2), [
    ("Step",  "dets (no trim)", "dets (TrimSQD)", "Energy (Ha)",      "Variance (Ha²)"),
    ("0",     "995",            "995",            "−101.9406",        "1.649"),
    ("3",     "11,042",         "5,794  (−47%)",  "−103.5938 (0.16 mHa from FCI)", "≤ 0.001"),
])

add_footer(slide, "Slide 8 — roadmap (co-presenter handles chemistry interpretation)")


# =============================================================================
# Slide 9 — The HPC pain (build matrix)
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "The HPC pain — toolchain × MPI × fabric per cluster",
          "Why standing this up on a new cluster is harder than pip install")

add_text_box(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4),
             ["Three backends, three compilers — distutils takes one CXX per setup() call → two pip invocations"],
             size=Pt(15), color=PRIMARY, bold_first=True)

add_table(slide, Inches(0.5), Inches(2.05), Inches(12.3), Inches(1.5), [
    ("Backend",                  "Compiler",                 "When"),
    ("_core_cpu",                "system c++",               "debugging, small problems"),
    ("_core_gpu (Thrust)",       "NVHPC nvc++",              "production GPU on NVIDIA"),
    ("_core_gpu_omp_nvidia",     "LLVM clang++ + NVPTX",     "alternative GPU path"),
])

add_text_box(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.4),
             ["Plus a long tail of environment plumbing per box:"],
             size=Pt(15), color=DARK, bold_first=True)

add_text_box(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8),
             ["•  GPU compute capability — sm_90 (default) vs sm_100 (Blackwell). Mismatch → silent host fallback, garbage energy.",
              "•  MPI vendor — HPCX 4.1.x can't init pml=ucx in SLURM cgroup; fall back to ob1+smcuda+tcp. OpenMPI 5.x is fine. Spectrum MPI needs jsrun.",
              "•  mpi4py ABI — wheels built against OpenMPI 5 segfault on HPCX 4. Source rebuild required.",
              "•  LD_LIBRARY_PATH ordering — LLVM's libomp must precede NVHPC's, or the wrong OpenMP runtime loads.",
              "•  SLURM DefMemPerCPU — opaque 'Requested node configuration is not available' if you forget --mem.",
              "•  cvd_wrapper local-rank source — SLURM_LOCALID is alloc-wide; OMPI_COMM_WORLD_LOCAL_RANK is post-MPI_Init; OMPI_MCA_orte_ess_node_rank is the right one.",
             ],
             size=Pt(13))

add_footer(slide, "Slide 9 — HPC pain (build matrix)")


# =============================================================================
# Slide 10 — Network fabrics matter
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Network fabric dominates at 2+ nodes",
          "Same SBD code; comm-backend choice changes wallclock by ~1.4×")

add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.5), [
    ("Cluster",                "Fabric",                          "Best comm",         "Why"),
    ("coreweave GB200 (2-node)","MNNVL across nodes",              "nccl (P2P/MNNVL)",   "NVLink cross-node, not IB"),
    ("8×H100 (1-node)",        "NVLink + NVSwitch + PCIe Gen5",   "nccl ≈ cuda_mpi",    "no MNNVL; comm dwarfed by compute"),
    ("IBM LSF / jsrun",        "Spectrum MPI + IB+GDR",           "cuda_mpi",           "IBM's CUDA-aware MPI"),
    ("Plain VM (no IB)",       "TCP only",                        "host_mpi or nccl",   "no high-speed cross-node fabric"),
])

add_text_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
             ["Concrete: same Fulqrum + SBD on 2-node GB200:",
              "  • nccl + MNNVL:    188 s",
              "  • cuda_mpi + TCP-staged IB:  264 s",
              "",
              "→ Fabric choice is a ~1.4× lever, independent of the kernel.",
              "→ Choosing wrong leaves perf on the table, silently."],
             size=Pt(15))

add_footer(slide, "Slide 10 — fabric matters")


# =============================================================================
# Slide 11 — The agent
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Letting an agent absorb the HPC stack",
          "UCX, NCCL, MPI, fabric, launcher — discovered per cluster, not encoded in your head")

add_text_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5),
             ["Single-node SBD is mostly tractable. Cluster-vendor stacks bite at 2+ nodes, where choices compound:"],
             size=Pt(15), color=DARK, bold_first=True)

add_text_box(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.5),
             ["•  UCX transport selection (cuda_copy / cuda_ipc / rc / ud) and which the cluster cgroup exposes",
              "•  NCCL version (≥ 2.23 for MNNVL) and which fabric it picks (P2P/MNNVL vs IB+GDR vs sockets)",
              "•  MPI vendor (HPCX vs OpenMPI 5 vs Spectrum) and the launcher (mpirun vs srun --mpi=pmix vs jsrun)",
              "•  GPU-pinning convention — which env var the wrapper reads to pick CUDA_VISIBLE_DEVICES per rank",
              "•  Fabric peculiarities (MNNVL on GB200, GPUDirect peermem on/off on H100, no IB on a VM)",
             ],
             size=Pt(14))

add_text_box(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0),
             ["The sqd-onboard agent walks the cluster, picks the right backend / launcher / fabric flags, builds the stack, validates correctness against a small reference, and emits ready-to-submit run scripts under run/<solver>/<n>node/.",
              "",
              "When it hits a failure mode it hasn't seen, the fix is captured as a signature in the playbook so the next user on that cluster doesn't repeat the discovery."],
             size=Pt(14), color=DIM)

add_footer(slide, "Slide 11 — the agent")


# =============================================================================
# Slide 12 — Backup: install
# =============================================================================
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Backup — install commands", None)

add_code_box(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.5), [
    "# 1. CPU + Thrust GPU (NVHPC nvc++ on PATH)",
    "SBD_BUILD_BACKEND=both CC=nvc CXX=nvc++ \\",
    "    pip install --no-build-isolation -e .",
    "",
    "# 2. OMP-offload GPU (LLVM clang++ with NVPTX target on PATH)",
    "#    Built separately — incompatible with NVHPC's CXX in one setup() call.",
    "SBD_BUILD_BACKEND=gpu_omp_nvidia \\",
    "    pip install --no-build-isolation -e .",
    "",
    "# 3. qiskit-addon-sqd MPI-aware fork",
    "pip install git+https://github.com/hfwen0502/qiskit-addon-sqd@patch-ferminon-sbd",
], highlight_idx=[1, 6, 10], size=Pt(14))

add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.0),
             ["Per-cluster envelope: this is what the agent (slide 11) handles for you.",
              "",
              "Reference perf: see .github/FULQRUM_SBD_GB200_SCALING.md for the full per-matvec breakdown."],
             size=Pt(15), color=DIM)

add_footer(slide, "Slide 12 — backup")


# -----------------------------------------------------------------------------
# Save
# -----------------------------------------------------------------------------
import sys
out = sys.argv[1] if len(sys.argv) > 1 else "/tmp/SBD_PYTHON_DECK.pptx"
prs.save(out)
print(f"wrote {out}")
